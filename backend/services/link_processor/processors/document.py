"""
services/link_processor/processors/document.py
- Reads and chunks PDF, DOCX, and TXT files.
"""

import logging
import asyncio
from pathlib import Path
from engine.utils import get_embedding as get_emb_util

logger = logging.getLogger(__name__)


def _read_pdf(file_path: str) -> str:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text.strip()
    except Exception as e:
        logger.error(f"PDF read error: {e}")
        return ""


def _read_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.error(f"DOCX read error: {e}")
        return ""


def _read_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.error(f"TXT read error: {e}")
        return ""


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list:
    """Split text into overlapping word-based chunks."""
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunks.append(" ".join(words[i:i + chunk_size]))
        i += chunk_size - overlap
    return chunks


def _read_csv(file_path: str) -> str:
    try:
        import polars as pl
        logger.info(f"⚡ [POLARS] Reading {file_path} extremely fast...")
        df = pl.read_csv(file_path, ignore_errors=True)
        sample = df.head(50)
        desc = df.describe()
        return f"POLARS EXTRACT: {len(df)} total rows.\n--- SCHEMA DEEP DIVE ---\n{desc}\n\n--- SAMPLE (TOP 50) ---\n{sample}"
    except Exception as e:
        logger.error(f"CSV read error: {e}")
        return ""


def _read_excel(file_path: str) -> str:
    try:
        import pandas as pd
        # Try to read all sheets
        dfs = pd.read_excel(file_path, sheet_name=None)
        all_text = []
        for sheet_name, df in dfs.items():
            all_text.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}")
        return "\n\n".join(all_text)
    except Exception as e:
        logger.error(f"Excel read error: {e}")
        return ""


def _read_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_runs = []
        for i, slide in enumerate(prs.slides):
            text_runs.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        logger.error(f"PPTX read error: {e}")
        return ""


async def process_document_file(file_path: str, filename: str) -> dict:
    """Read and chunk a document file. Returns structured content with embeddings."""
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        text = _read_pdf(file_path)
        fmt = "PDF"
    elif ext in (".docx", ".doc"):
        text = _read_docx(file_path)
        fmt = "DOCX"
    elif ext == ".csv":
        text = _read_csv(file_path)
        fmt = "CSV"
    elif ext in (".xlsx", ".xls"):
        text = _read_excel(file_path)
        fmt = "EXCEL"
    elif ext == ".pptx":
        text = _read_pptx(file_path)
        fmt = "POWERPOINT"
    else:
        text = _read_txt(file_path)
        fmt = "TXT"

    if not text:
        return {"error": "Could not extract text", "text": "", "chunks": [], "embeddings": [], "source": filename}

    chunks = chunk_text(text)
    
    # Generate embeddings for chunks in parallel
    embeddings = []
    try:
        tasks = [get_emb_util(chunk) for chunk in chunks[:100]] # Increased to 100 for better coverage
        embeddings = await asyncio.gather(*tasks)
    except Exception as e:
        logger.error(f"Error generating chunk embeddings: {e}")

    return {
        "source": filename,
        "title": filename,
        "format": fmt,
        "text": text[:5000],           # First 5000 chars for summary
        "full_text": text,
        "chunks": chunks,
        "embeddings": embeddings,
        "word_count": len(text.split()),
        "error": None,
    }
